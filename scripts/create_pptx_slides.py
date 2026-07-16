"""Generate PPTX presentation from HUST template with Graph ML project content.

Usage:
    python scripts/create_pptx_slides.py
"""
from __future__ import annotations

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# --- Paths ---
BRANDING = "/Users/ducnguyen/Projects/HUST2/IT5315E/reports/branding"
TEMPLATE = os.path.join(BRANDING, "hust-template-4x3.pptx")
LOGO = os.path.join(BRANDING, "hust-logo.png")
EMBLEM = os.path.join(BRANDING, "hust-emblem.png")
PROJECT = "/Users/ducnguyen/Projects/HUST2/IT5429E/GraphMl"
FIGURES = os.path.join(PROJECT, "figures")
OUTPUT = os.path.join(PROJECT, "report", "graphml-presentation.pptx")

# --- Colors ---
DARK = RGBColor(0x1A, 0x1A, 0x2E)
PRIMARY = RGBColor(0x00, 0x6B, 0x5A)  # HUST green-ish
ACCENT = RGBColor(0xC0, 0x39, 0x2B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0x99, 0x99, 0x99)


def set_text(tf, text, size=14, bold=False, color=None, alignment=PP_ALIGN.LEFT):
    """Set text frame with a single paragraph."""
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def add_bullet_paragraph(tf, text, size=13, bold=False, color=None, indent=0):
    """Add a bullet point to an existing text frame."""
    p = tf.add_paragraph()
    p.space_before = Pt(4)
    p.space_after = Pt(2)
    p.level = indent
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return p


def add_slide(prs, layout_idx):
    """Add a slide with the given layout index."""
    return prs.slides.add_slide(prs.slide_layouts[layout_idx])


def set_title(slide, text, size=24):
    """Set title text on a slide (placeholder 0)."""
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 0:
            set_text(shape.text_frame, text, size=size, bold=True, color=DARK)
            return
    # Fallback: try Title 1 or Title 2
    for shape in slide.shapes:
        if "Title" in shape.name and hasattr(shape, "text_frame"):
            set_text(shape.text_frame, text, size=size, bold=True, color=DARK)
            return


def add_textbox(slide, left, top, width, height):
    """Add a text box and return it."""
    return slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))


def main():
    prs = Presentation(TEMPLATE)

    # Remove all template slides (keep layouts)
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

    # ============================================================
    # SLIDE 1: Title
    # ============================================================
    slide = add_slide(prs, 1)  # Title Slide layout
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 0:
            tf = shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = "Residual Multi-Head Mixing GAT"
            run.font.size = Pt(28)
            run.font.bold = True
            run.font.color.rgb = DARK

            add_bullet_paragraph(tf, "Semi-Supervised Node Classification on Cora",
                                 size=18, color=GRAY)
            add_bullet_paragraph(tf, "", size=10)
            add_bullet_paragraph(tf, "Graph ML Project — Hanoi University of Science and Technology",
                                 size=12, color=LIGHT_GRAY)
            add_bullet_paragraph(tf, "", size=8)
            add_bullet_paragraph(tf,
                "Do Tuan Minh · Tran Tien Dung · Tran Manh Tien · Hoang Huy Chien · Nguyen Tien Duc",
                size=11, color=GRAY)

    # ============================================================
    # SLIDE 2: Problem Statement
    # ============================================================
    slide = add_slide(prs, 4)  # Two Content layout
    set_title(slide, "Problem: Semi-Supervised Node Classification")

    # Left column: problem
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            tf = shape.text_frame
            tf.clear()
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "Problem"
            run.font.size = Pt(16)
            run.font.bold = True
            run.font.color.rgb = PRIMARY

            for item in [
                "Given graph G = (V, E) — Cora citation network",
                "2708 nodes (papers) · 5278 edges (citations)",
                "Each node: 1433-dim bag-of-words features",
                "Only 140 nodes have training labels (5.2%)",
                "Classify into 7 classes (research topics)",
            ]:
                add_bullet_paragraph(tf, f"• {item}", size=12)

        # Right column: why GNN
        elif shape.placeholder_format.idx == 2:
            tf = shape.text_frame
            tf.clear()
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "Why Graph Neural Networks?"
            run.font.size = Pt(16)
            run.font.bold = True
            run.font.color.rgb = PRIMARY

            for item in [
                "Citation structure carries important information",
                "Papers in the same field tend to cite each other",
                "MLP ignores graph structure → loses information",
                "GNN: message passing between neighbors",
                "Expectation: GNN >> MLP on this task",
            ]:
                add_bullet_paragraph(tf, f"• {item}", size=12)

    # ============================================================
    # SLIDE 3: Dataset
    # ============================================================
    slide = add_slide(prs, 5)  # Comparison layout
    set_title(slide, "Dataset: Cora Citation Network")

    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 13:  # Text Placeholder
            tf = shape.text_frame
            tf.clear()
            tf.word_wrap = True

            stats = [
                ("Nodes (papers)", "2,708"),
                ("Edges (citations)", "5,278 (directed) → 10,556 (undirected)"),
                ("Features", "1,433 (bag-of-words, binary)"),
                ("Classes", "7 (research topics)"),
                ("Train / Val / Test", "140 / 500 / 1,000"),
                ("Sparsity", "~0.14% (very sparse graph)"),
            ]

            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "Cora Statistics"
            run.font.size = Pt(16)
            run.font.bold = True
            run.font.color.rgb = PRIMARY

            for label, val in stats:
                add_bullet_paragraph(tf, f"• {label}: {val}", size=13)

            add_bullet_paragraph(tf, "", size=8)
            add_bullet_paragraph(tf, "Data Preparation:", size=14, bold=True, color=PRIMARY)
            for item in [
                "Planetoid loader (PyG) — standard Cora split",
                "NormalizeFeatures transform — row-normalize",
                "Fixed train/val/test mask (reproducible)",
            ]:
                add_bullet_paragraph(tf, f"• {item}", size=12)

    # ============================================================
    # SLIDE 4: Baseline Models
    # ============================================================
    slide = add_slide(prs, 5)  # Comparison layout
    set_title(slide, "Baseline Models")

    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 13:
            tf = shape.text_frame
            tf.clear()
            tf.word_wrap = True

            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "4 Baseline Models"
            run.font.size = Pt(16)
            run.font.bold = True
            run.font.color.rgb = PRIMARY

            models = [
                ("MLP", "2-layer fully connected, ignores graph structure. Weakest baseline."),
                ("GCN", "Graph Convolutional Network. Spectral convolution with normalized adjacency matrix."),
                ("GraphSAGE", "SAmple and aggreGatE. Inductive learning with mean aggregator."),
                ("GAT", "Graph Attention Network. Learned attention weights, multi-head attention."),
            ]
            for name, desc in models:
                add_bullet_paragraph(tf, "", size=4)
                add_bullet_paragraph(tf, f"▸ {name}", size=14, bold=True, color=DARK)
                add_bullet_paragraph(tf, f"  {desc}", size=12, color=GRAY)

    # ============================================================
    # SLIDE 5: Proposed Model
    # ============================================================
    slide = add_slide(prs, 5)  # Comparison layout
    set_title(slide, "Proposed Model: Residual Multi-Head Mixing GAT")

    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 13:
            tf = shape.text_frame
            tf.clear()
            tf.word_wrap = True

            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "Architecture"
            run.font.size = Pt(16)
            run.font.bold = True
            run.font.color.rgb = PRIMARY

            lines = [
                "H₀ = X                                    # Raw features (1433-dim)",
                "H₁ = ELU( GATv2_multi(H₀, E) )     # 8 heads, concat → 64-dim",
                "H₂ = GATv2_single(H₁, E)               # Single head → 7-dim",
                "Z  = H₂ + H₀ · W_skip              # Residual skip → logits",
            ]
            for line in lines:
                add_bullet_paragraph(tf, f"  {line}", size=11)

            add_bullet_paragraph(tf, "", size=6)
            add_bullet_paragraph(tf, "3 Key Contributions:", size=14, bold=True, color=PRIMARY)

            contribs = [
                ("GATv2 (replacing GAT)", "Dynamic attention coefficients, avoids static ranking problem"),
                ("Residual Skip Connection", "H0·W_skip: direct gradient flow, counteracts over-smoothing"),
                ("DropEdge (p=0.2)", "Random edge removal during training → structural augmentation"),
            ]
            for name, desc in contribs:
                add_bullet_paragraph(tf, f"  ▸ {name}", size=13, bold=True)
                add_bullet_paragraph(tf, f"    {desc}", size=11, color=GRAY)

    # ============================================================
    # SLIDE 6: Experiment Setup
    # ============================================================
    slide = add_slide(prs, 4)  # Two Content layout
    set_title(slide, "Experimental Setup")

    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            tf = shape.text_frame
            tf.clear()
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "Hyperparameters"
            run.font.size = Pt(16)
            run.font.bold = True
            run.font.color.rgb = PRIMARY

            params = [
                "Optimizer: Adam",
                "Learning rate: 0.005",
                "Weight decay: 5×10⁻⁴",
                "Epochs: 200",
                "Early stopping: patience=100",
                "Dropout: 0.6",
                "Hidden dim (GCN/SAGE): 64",
                "GAT heads: 8 (per-head=8)",
            ]
            for item in params:
                add_bullet_paragraph(tf, f"• {item}", size=12)

        elif shape.placeholder_format.idx == 2:
            tf = shape.text_frame
            tf.clear()
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "Protocol"
            run.font.size = Pt(16)
            run.font.bold = True
            run.font.color.rgb = PRIMARY

            proto = [
                "5 random seeds (42-46)",
                "Report: mean ± std",
                "Model selection: best val accuracy",
                "Report: test acc at best-val epoch",
                "Standard Cora split: 140/500/1000",
                "NormalizeFeatures transform",
                "PyTorch 2.13 + PyG 2.8",
                "Full reproducibility: seed all RNGs",
            ]
            for item in proto:
                add_bullet_paragraph(tf, f"• {item}", size=12)

    # ============================================================
    # SLIDE 7: Results — Model Comparison (with chart image)
    # ============================================================
    slide = add_slide(prs, 4)  # Two Content layout
    set_title(slide, "Results: Model Comparison")

    # Left: bar chart image
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            img_path = os.path.join(FIGURES, "accuracy_bar_chart.png")
            if os.path.exists(img_path):
                slide.shapes.add_picture(
                    img_path,
                    shape.left, shape.top,
                    shape.width, shape.height
                )

        # Right: table
        elif shape.placeholder_format.idx == 2:
            tf = shape.text_frame
            tf.clear()
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "Results (5 seeds, 200 epochs)"
            run.font.size = Pt(14)
            run.font.bold = True
            run.font.color.rgb = PRIMARY

            results = [
                ("MLP", "0.5888 ± 0.0104"),
                ("GraphSAGE", "0.8130 ± 0.0077"),
                ("GAT", "0.8272 ± 0.0080"),
                ("GCN", "0.8276 ± 0.0050"),
                ("Proposed ★", "0.8326 ± 0.0080"),
            ]
            for name, acc in results:
                is_best = "★" in name
                add_bullet_paragraph(tf, f"• {name}: {acc}",
                                     size=13, bold=is_best,
                                     color=PRIMARY if is_best else DARK)

            add_bullet_paragraph(tf, "", size=8)
            add_bullet_paragraph(tf, "Key Takeaways:", size=13, bold=True, color=PRIMARY)
            add_bullet_paragraph(tf, "• MLP → GNN: +24% (graph structure is crucial)", size=11)
            add_bullet_paragraph(tf, "• Proposed outperforms GAT by +0.5%", size=11)
            add_bullet_paragraph(tf, "• Low std → stable training", size=11)

    # ============================================================
    # SLIDE 8: Training Curves
    # ============================================================
    slide = add_slide(prs, 7)  # Title Only layout
    set_title(slide, "Training Curves: Loss & Validation Accuracy")

    # Add the curves image centered
    img_path = os.path.join(FIGURES, "model_comparison_curves.png")
    if os.path.exists(img_path):
        slide.shapes.add_picture(
            img_path,
            Inches(0.5), Inches(1.2),
            Inches(9.0), Inches(5.5)
        )

    # ============================================================
    # SLIDE 9: Ablation Study
    # ============================================================
    slide = add_slide(prs, 4)  # Two Content layout
    set_title(slide, "Ablation Study: Residual Multi-Head Mixing GAT")

    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            img_path = os.path.join(FIGURES, "ablation_bar_chart.png")
            if os.path.exists(img_path):
                slide.shapes.add_picture(
                    img_path,
                    shape.left, shape.top,
                    shape.width, shape.height
                )

        elif shape.placeholder_format.idx == 2:
            tf = shape.text_frame
            tf.clear()
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "Component Analysis"
            run.font.size = Pt(14)
            run.font.bold = True
            run.font.color.rgb = PRIMARY

            ablations = [
                ("GAT (base)", "0.8164", "—"),
                ("+ residual", "0.8334", "+1.7%  ✓"),
                ("+ res + dropedge", "0.8326", "+1.6%"),
                ("full, heads=4", "0.8260", "+1.0%"),
            ]
            for name, acc, delta in ablations:
                add_bullet_paragraph(tf, f"• {name}: {acc}  ({delta})", size=12)

            add_bullet_paragraph(tf, "", size=6)
            add_bullet_paragraph(tf, "Findings:", size=13, bold=True, color=PRIMARY)
            add_bullet_paragraph(tf, "✓ Residual skip: +1.7% — main contribution", size=11, color=PRIMARY)
            add_bullet_paragraph(tf, "≈ DropEdge: −0.1% — negligible on small graph", size=11, color=GRAY)
            add_bullet_paragraph(tf, "↓ Heads 8→4: −0.7% — more heads help", size=11, color=ACCENT)

    # ============================================================
    # SLIDE 10: Confusion Matrix
    # ============================================================
    slide = add_slide(prs, 7)  # Title Only layout
    set_title(slide, "Confusion Matrix: Proposed Model on Test Set")

    img_path = os.path.join(FIGURES, "confusion_matrix.png")
    if os.path.exists(img_path):
        slide.shapes.add_picture(
            img_path,
            Inches(1.5), Inches(1.2),
            Inches(7.0), Inches(5.8)
        )

    # ============================================================
    # SLIDE 11: Discussion
    # ============================================================
    slide = add_slide(prs, 4)  # Two Content layout
    set_title(slide, "Discussion")

    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            tf = shape.text_frame
            tf.clear()
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "Key Findings"
            run.font.size = Pt(16)
            run.font.bold = True
            run.font.color.rgb = PRIMARY

            findings = [
                "Graph structure is crucial: MLP 0.59 → GNN 0.82 (+24%)",
                "Residual skip connection works: +1.7% test accuracy",
                "Stable training: low std (0.005-0.014)",
                "Val-test gap ~1-2%: no severe overfitting",
            ]
            for item in findings:
                add_bullet_paragraph(tf, f"✓ {item}", size=12, color=PRIMARY)

        elif shape.placeholder_format.idx == 2:
            tf = shape.text_frame
            tf.clear()
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "Surprises"
            run.font.size = Pt(16)
            run.font.bold = True
            run.font.color.rgb = ACCENT

            surprises = [
                ("GCN ≈ GAT", "Attention does not outperform simple mean aggregation on Cora"),
                ("DropEdge ≈ 0", "Regularization unnecessary on such a small graph"),
            ]
            for title, desc in surprises:
                add_bullet_paragraph(tf, f"✗ {title}", size=13, bold=True, color=ACCENT)
                add_bullet_paragraph(tf, f"  {desc}", size=11, color=GRAY)

            add_bullet_paragraph(tf, "", size=8)
            add_bullet_paragraph(tf, "Explanation:", size=13, bold=True, color=DARK)
            add_bullet_paragraph(tf, "Cora is small (2708 nodes) and homogeneous → simple GNNs are already sufficient. DropEdge and attention shine on larger graphs.", size=11, color=GRAY)

    # ============================================================
    # SLIDE 12: Limitations & Future Work
    # ============================================================
    slide = add_slide(prs, 4)  # Two Content layout
    set_title(slide, "Limitations & Future Work")

    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:
            tf = shape.text_frame
            tf.clear()
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "Limitations"
            run.font.size = Pt(16)
            run.font.bold = True
            run.font.color.rgb = ACCENT

            limitations = [
                "Single dataset only (Cora, 2708 nodes)",
                "Shallow architecture (2 layers)",
                "No systematic hyperparameter tuning",
                "Accuracy metric only (no F1, AUC)",
                "Limited regularization study (no Mixup, Label Smoothing)",
            ]
            for item in limitations:
                add_bullet_paragraph(tf, f"• {item}", size=12)

        elif shape.placeholder_format.idx == 2:
            tf = shape.text_frame
            tf.clear()
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "Future Work"
            run.font.size = Pt(16)
            run.font.bold = True
            run.font.color.rgb = PRIMARY

            future = [
                "Citeseer, Pubmed, Arxiv, OGB datasets",
                "Deeper GNN (3-5 layers) + Jumping Knowledge",
                "Graph Transformer (global attention)",
                "Attention weight visualization",
                "Real-world apps: recommendation, social networks",
            ]
            for item in future:
                add_bullet_paragraph(tf, f"→ {item}", size=12, color=PRIMARY)

    # ============================================================
    # SLIDE 13: Conclusion & Q&A
    # ============================================================
    slide = add_slide(prs, 5)  # Comparison layout
    set_title(slide, "Conclusion")

    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 13:
            tf = shape.text_frame
            tf.clear()
            tf.word_wrap = True

            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = "Summary"
            run.font.size = Pt(18)
            run.font.bold = True
            run.font.color.rgb = PRIMARY

            add_bullet_paragraph(tf, "", size=8)

            conclusions = [
                "0.8326 — Best test accuracy (Proposed model)",
                "+1.7% — Residual skip connection improvement",
                "+24% — Graph structure vs MLP improvement",
            ]
            for item in conclusions:
                p = add_bullet_paragraph(tf, item, size=16, bold=True, color=DARK)
                p.alignment = PP_ALIGN.CENTER

            add_bullet_paragraph(tf, "", size=16)

            p = add_bullet_paragraph(tf,
                "Residual Multi-Head Mixing GAT successfully combines\n"
                "multi-head attention with an input-feature skip connection.",
                size=14, color=GRAY)
            p.alignment = PP_ALIGN.CENTER

            add_bullet_paragraph(tf, "", size=24)

            p = add_bullet_paragraph(tf, "Thank You & Q&A", size=22, bold=True, color=PRIMARY)
            p.alignment = PP_ALIGN.CENTER

            p = add_bullet_paragraph(tf, "Code: github.com/minhdo2207/GraphMl", size=11, color=LIGHT_GRAY)
            p.alignment = PP_ALIGN.CENTER

    # --- Save ---
    os.makedirs(os.path.dirname(OUTPUT), exist_ok=True)
    prs.save(OUTPUT)
    print(f"Saved: {OUTPUT}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
